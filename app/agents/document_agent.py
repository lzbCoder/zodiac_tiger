import io
import re
from datetime import datetime
from pathlib import Path

from langchain_core.runnables import RunnableConfig
from langchain_core.callbacks import adispatch_custom_event
from loguru import logger

from app.state.agent_state import AgentState
from app.services import file_service


async def _stream_doc_info(text: str, config: RunnableConfig) -> None:
    """把确定性文件信息按行作为 doc_token 自定义事件发出，经 event_sse 转 token 上屏。

    不经 LLM，保证下载链接等内容逐字正确。失败不阻断主流程。
    """
    try:
        lines = text.split("\n")
        for i, line in enumerate(lines):
            chunk = line if i == len(lines) - 1 else line + "\n"
            await adispatch_custom_event("doc_token", {"content": chunk}, config=config)
    except Exception as e:
        logger.warning(f"[文档] 流式输出文件信息失败: {e}")

FORMAT_EXT: dict[str, str] = {
    "md": ".md",
    "docx": ".docx",
    "xlsx": ".xlsx",
    "html": ".html",
    "ppt": ".pptx",
    "pptx": ".pptx",
}


def _suggest_filename(content: str, ext: str) -> str:
    """从内容首行提取合适的文件名。"""
    lines = content.strip().split("\n")
    base = "document"
    for line in lines:
        clean = re.sub(r"^#+\s*", "", line).strip()
        if not clean:
            continue
        clean = re.sub(r"[*_~`|\\/:<>?\"']", "", clean)
        clean = re.sub(r"\s+", " ", clean).strip()
        if clean:
            base = clean[:10]
            break
    return f"{base}{ext}"


def _generate_md_bytes(content: str) -> bytes:
    return content.encode("utf-8")


def _generate_docx_bytes(content: str) -> bytes:
    from docx import Document
    doc = Document()
    for line in content.split("\n"):
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("- ") or line.startswith("* "):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif line.strip():
            doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _generate_xlsx_bytes(content: str) -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    for row_idx, line in enumerate(content.split("\n"), start=1):
        parts = [c.strip() for c in line.split("|") if c.strip()]
        for col_idx, val in enumerate(parts, start=1):
            ws.cell(row=row_idx, column=col_idx, value=val)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _generate_html_bytes(content: str) -> bytes:
    paragraphs = []
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            paragraphs.append(f"<h1>{line[2:]}</h1>")
        elif line.startswith("## "):
            paragraphs.append(f"<h2>{line[3:]}</h2>")
        elif line.startswith("### "):
            paragraphs.append(f"<h3>{line[4:]}</h3>")
        elif line.startswith("- "):
            paragraphs.append(f"<li>{line[2:]}</li>")
        else:
            paragraphs.append(f"<p>{line}</p>")
    body = "\n".join(paragraphs)

    html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head><meta charset="UTF-8"><title>Document</title>
<style>
  body {{ font-family: 'Microsoft YaHei', sans-serif; max-width: 900px; margin: 40px auto; padding: 0 20px; color: #333; line-height: 1.8; }}
  h1 {{ border-bottom: 2px solid #007acc; padding-bottom: 8px; }}
  h2 {{ color: #007acc; }}
  table {{ border-collapse: collapse; width: 100%; }}
  th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
  th {{ background: #007acc; color: white; }}
  code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 3px; }}
  pre {{ background: #2d2d2d; color: #eee; padding: 16px; border-radius: 6px; overflow-x: auto; }}
</style></head>
<body>{body}</body></html>"""
    return html.encode("utf-8")


def _generate_pptx_bytes(content: str) -> bytes:
    """按 markdown 标题切分幻灯片，标题下的文本/列表作为正文要点。"""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_layout = prs.slide_layouts[0]      # 标题幻灯片
    content_layout = prs.slide_layouts[1]    # 标题 + 内容

    slides: list[tuple[str, list[str]]] = []
    cur_title = "演示文稿"
    cur_bullets: list[str] = []
    has_slide = False

    def _flush() -> None:
        nonlocal cur_bullets
        if has_slide or cur_bullets:
            slides.append((cur_title, cur_bullets))
        cur_bullets = []

    for raw in content.split("\n"):
        line = raw.rstrip()
        if not line.strip():
            continue
        if line.startswith("# ") or line.startswith("## ") or line.startswith("### "):
            _flush()
            cur_title = re.sub(r"^#+\s*", "", line).strip()
            has_slide = True
        else:
            bullet = re.sub(r"^[-*]\s+", "", line).strip()
            if bullet:
                cur_bullets.append(bullet)
    _flush()

    if not slides:
        slides = [("演示文稿", [content.strip()[:200]])]

    # 首页：标题版式，用第一张幻灯片的标题作为封面标题
    cover = prs.slides.add_slide(title_layout)
    cover.shapes.title.text = slides[0][0]

    # 其余每张：标题 + 内容版式，要点逐条成段
    for title, bullets in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = b
            para.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_FILE_GENERATORS = {
    "md": _generate_md_bytes,
    "docx": _generate_docx_bytes,
    "xlsx": _generate_xlsx_bytes,
    "html": _generate_html_bytes,
    "ppt": _generate_pptx_bytes,
    "pptx": _generate_pptx_bytes,
}


async def _generate_file_bytes(content: str, fmt: str) -> tuple[bytes, str]:
    """生成文件字节到内存，不写盘。"""
    ext = FORMAT_EXT.get(fmt, ".md")
    file_name = _suggest_filename(content, ext)
    generator = _FILE_GENERATORS.get(fmt, _generate_md_bytes)
    file_bytes = generator(content)
    logger.info(f"文件内容已生成: {file_name} ({fmt}, {len(file_bytes)} bytes)")
    return file_bytes, file_name


async def document_agent_node(state: AgentState, config: RunnableConfig) -> dict:
    content = state.get("generate_content", "")
    fmt = state.get("generate_format", "").lower().strip()
    session_id = config["configurable"]["session_id"]
    chat_id = config["configurable"]["chat_id"]

    if not content or not fmt or fmt == "none" or fmt == "md":
        return {}

    # ---- 不支持的类型 → 降级提示 + 输出 MD 原文 ----
    if fmt not in FORMAT_EXT:
        hint = (
            f"暂不支持生成 **{fmt.upper()}** 格式文档，"
            f"已为您默认输出 Markdown 原文：\n\n{content}"
        )
        await _stream_doc_info(hint, config)
        return {"messages": [{"role": "ai", "content": hint}]}

    # ---- 实体文件生成（内存 → save_file 写盘 + 入库） ----
    try:
        file_bytes, file_name = await _generate_file_bytes(content, fmt)
    except Exception as e:
        logger.error(f"文件生成失败: {e}")
        return {}

    try:
        record = await file_service.save_file(
            session_id=session_id,
            chat_id=chat_id,
            file_name=file_name,
            content=file_bytes,
            created_by="文档生成 Agent",
        )
        file_size_str = file_service.format_size(record.file_size)
        created_at_str = record.created_at.strftime("%Y-%m-%d %H:%M:%S") if record.created_at else ""
        download_url = f"/api/file/download/{record.id}"
    except Exception as e:
        logger.error(f"文件入库失败: {e}")
        file_size_str = file_service.format_size(len(file_bytes))
        created_at_str = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        download_url = "#"

    reply = (
        f"文档已生成完毕！\n\n"
        f"**【文件信息】**\n"
        f"- 文件名：{file_name}\n"
        f"- 存储路径：files/{datetime.now().strftime('%Y-%m-%d')}/{session_id}/{file_name}\n"
        f"- 文件大小：{file_size_str}\n"
        f"- 创建时间：{created_at_str}\n\n"
        f"[点击下载文档]({download_url})"
    )

    await _stream_doc_info(reply, config)   # 确定性文件信息流式上屏（不经 LLM）

    result: dict = {"messages": [{"role": "ai", "content": reply}]}
    try:
        result["last_file_id"] = str(record.id)   # 供 artifact_store 关联产物文件
    except Exception:
        pass
    return result
